import os

def extract_text_from_file(file_path: str) -> str:
    """
    Universally detects and extracts readable text from almost any file format.
    Features robust error handling to prevent 500 Server Crashes.
    """
    # Ensure the file actually exists before we try to read it
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Cannot find the file to parse: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        # 1. Plain Text, Code, Logs, and Scripts
        if ext in ['.txt', '.md', '.csv', '.json', '.html', '.py', '.js', '.http', '.log', '.ini', '.cfg', '.env']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                # Fallback for weirdly encoded text files
                with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
                    return f.read()
                
        # 2. PDF Documents
        elif ext == '.pdf':
            try:
                import pymupdf4llm
                text = pymupdf4llm.to_markdown(file_path)
                if not text.strip():
                    raise ValueError("Extracted PDF text is empty.")
                return text
            except Exception as e:
                # Secondary Fallback for PDFs if the Markdown converter fails
                import fitz
                doc = fitz.open(file_path)
                
                page_texts = []
                for page in doc:
                    # Defensive Check: Ensure the page object exists and has the extraction method
                    if page is not None and hasattr(page, 'get_text'):
                        page_texts.append(page.get_text())
                    else:
                        page_texts.append("[SYSTEM NOTE: A malformed or unreadable page element was skipped here.]")
                        
                return "\n".join(page_texts)
            
        # 3. Microsoft Word Documents (.docx)
        elif ext in ['.docx']:
            from docx import Document
            doc = Document(file_path)
            extracted = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            if not extracted:
                return "[SYSTEM NOTE: The uploaded Word Document appears to be empty or contains only images.]"
            return extracted
            
        # 4. Microsoft Excel Spreadsheets (.xlsx, .xls)
        elif ext in ['.xlsx', '.xls']:
            import pandas as pd
            # openpyxl is required to read modern Excel files
            df = pd.read_excel(file_path, engine='openpyxl')
            return df.to_markdown(index=False)
            
        # 5. Microsoft PowerPoint Presentations (.pptx)
        elif ext in ['.pptx']:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
            extracted = "\n\n".join(text_runs)
            if not extracted:
                return "[SYSTEM NOTE: The uploaded PowerPoint appears to be empty or contains only images.]"
            return extracted
            
        # 6. Local Images (.png, .jpg, .jpeg)
        elif ext in ['.png', '.jpg', '.jpeg']:
            try:
                from PIL import Image
                import pytesseract
                # Attempt to OCR the image
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                if not text.strip():
                    return "[SYSTEM NOTE: The system could not detect any readable text in this image.]"
                return text
            except Exception as img_e:
                return f"[SYSTEM NOTE: Image processing failed. Tesseract-OCR may not be installed on the host machine. Error: {str(img_e)}]"
            
        # 7. THE UNIVERSAL FALLBACK (Prevents 500 Crashes for Unknown Formats)
        else:
            try:
                # Try to forcefully extract any human-readable strings from the binary junk
                with open(file_path, 'rb') as f:
                    content = f.read()
                    # Filter out non-printable binary characters
                    readable = ''.join(chr(byte) for byte in content if 32 <= byte <= 126 or byte in (9, 10, 13))
                    if len(readable.strip()) > 50:
                        return f"[SYSTEM NOTE: Unknown format '{ext}'. Attempted forced text extraction.]\n\n" + readable
                    else:
                        return f"[SYSTEM NOTE: The file '{ext}' is a binary format that CipherMind cannot read.]"
            except Exception:
                return f"[SYSTEM NOTE: CipherMind does not currently support '{ext}' files.]"
                
    except ImportError as ie:
        return f"[SYSTEM ERROR: The backend is missing the required library to read {ext} files. Please install it: {str(ie)}]"
    except Exception as e:
        # We return the error as text instead of raising it, so the AI can explain the failure to the user instead of crashing the server!
        return f"[SYSTEM FATAL ERROR: The document parser encountered a critical failure while reading the {ext} file: {str(e)}]"